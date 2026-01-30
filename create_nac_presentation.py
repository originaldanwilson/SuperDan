"""
PowerPoint Presentation Generator for Network as Code (NaC) Transformation
Creates comprehensive deck explaining the migration from CLI/GUI to DevOps approach
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_title_slide(prs, title, subtitle):
    """Create title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide

def create_content_slide(prs, title, content_points):
    """Create bullet point content slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    
    text_frame = slide.placeholders[1].text_frame
    text_frame.clear()
    
    for point in content_points:
        p = text_frame.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(18)
    
    return slide

def create_two_column_slide(prs, title, left_content, right_content):
    """Create two-column comparison slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4), Inches(5))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    for item in left_content:
        p = left_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(5), Inches(1.5), Inches(4), Inches(5))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    for item in right_content:
        p = right_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
    
    return slide

def main():
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    create_title_slide(prs, 
                      "Network as Code (NaC) Transformation",
                      "Modernizing Network Operations with DevOps Practices\nGitHub + Ansible + Cisco NaC Collections")
    
    # Slide 2: Executive Summary
    create_content_slide(prs, "Executive Summary", [
        "Transform network operations from manual CLI/GUI to automated DevOps approach",
        "Leverage GitHub for version control and collaboration",
        "Use GitHub Actions for CI/CD automation",
        "Execute changes via Ansible with Cisco NaC collections",
        "Benefits: Speed, reliability, auditability, and scalability"
    ])
    
    # Slide 3: The Problem - Current State
    create_content_slide(prs, "Current State Challenges", [
        "❌ Manual CLI/GUI operations are time-consuming and error-prone",
        "❌ No version control or change tracking",
        "❌ Difficult to audit who made what changes and when",
        "❌ Configuration drift across devices",
        "❌ Limited ability to scale operations",
        "❌ Knowledge siloed in individual team members",
        "❌ No standardization or testing before deployment"
    ])
    
    # Slide 4: Why Network as Code?
    create_content_slide(prs, "Why Network as Code?", [
        "✅ Infrastructure as Code principles applied to networking",
        "✅ Version control for all network configurations (Git)",
        "✅ Automated testing and validation before deployment",
        "✅ Complete audit trail and rollback capabilities",
        "✅ Faster change implementation with lower risk",
        "✅ Standardized, repeatable processes",
        "✅ Self-documenting configurations",
        "✅ Collaboration through pull requests and code reviews"
    ])
    
    # Slide 5: Before vs After Comparison
    left = [
        "BEFORE (CLI/GUI)",
        "",
        "• Manual login to devices",
        "• Copy-paste configurations",
        "• No peer review",
        "• Manual documentation",
        "• Hard to replicate",
        "• No automated testing",
        "• Errors discovered in production",
        "• Difficult rollback"
    ]
    
    right = [
        "AFTER (NaC)",
        "",
        "• Define changes in YAML",
        "• Commit to GitHub",
        "• Automated peer review",
        "• Self-documenting code",
        "• Reusable templates",
        "• Automated validation",
        "• Errors caught pre-deployment",
        "• One-click rollback"
    ]
    
    create_two_column_slide(prs, "Transformation: Before vs After", left, right)
    
    # Slide 6: Solution Architecture
    create_content_slide(prs, "Solution Architecture Overview", [
        "1️⃣ GitHub: Central repository for YAML configurations",
        "2️⃣ GitHub Actions: CI/CD pipeline automation",
        "3️⃣ Ansible: Configuration management engine",
        "4️⃣ Cisco NaC Ansible Collections: Device-specific modules",
        "5️⃣ Network Devices: Cisco routers, switches, firewalls"
    ])
    
    # Slide 7: How It Works - The Workflow
    create_content_slide(prs, "How It Works: The NaC Workflow", [
        "Step 1: Engineer creates/modifies YAML configuration file",
        "Step 2: Commit and push changes to GitHub branch",
        "Step 3: Create Pull Request for peer review",
        "Step 4: GitHub Actions automatically validates YAML syntax",
        "Step 5: Team reviews and approves changes",
        "Step 6: Merge triggers GitHub Actions workflow",
        "Step 7: Ansible playbook executes against target devices",
        "Step 8: Changes applied, results logged and reported"
    ])
    
    # Slide 8: Technology Stack
    create_content_slide(prs, "Technology Stack", [
        "📦 GitHub Enterprise/Cloud: Version control & collaboration",
        "⚙️ GitHub Actions: CI/CD automation platform",
        "🤖 Ansible: Open-source automation engine",
        "📚 Cisco NaC Collections: cisco.ios, cisco.nxos, cisco.asa, etc.",
        "📝 YAML: Human-readable configuration format",
        "🔐 GitHub Secrets: Secure credential management",
        "📊 Ansible Tower/AWX (Optional): Enterprise automation platform"
    ])
    
    # Slide 9: Key Benefits
    create_content_slide(prs, "Key Benefits", [
        "🚀 Speed: Deploy changes in minutes, not hours",
        "🎯 Accuracy: Eliminate human error through automation",
        "📋 Compliance: Complete audit trail for all changes",
        "🔄 Repeatability: Consistent configurations across devices",
        "👥 Collaboration: Team-based change management",
        "🧪 Testing: Validate before production deployment",
        "⏮️ Rollback: Quick recovery from issues",
        "📈 Scalability: Manage thousands of devices efficiently"
    ])
    
    # Slide 10: Example YAML Configuration
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Example: YAML Configuration"
    
    yaml_example = """# Configure VLAN on switch
---
- name: Configure VLANs
  hosts: datacenter_switches
  gather_facts: no
  tasks:
    - name: Create VLAN 100
      cisco.ios.ios_vlans:
        config:
          - vlan_id: 100
            name: ENGINEERING
            state: active
        state: merged"""
    
    text_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = yaml_example
    p.font.name = 'Courier New'
    p.font.size = Pt(14)
    
    # Slide 11: GitHub Actions Workflow
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "GitHub Actions Workflow Example"
    
    workflow_example = """# .github/workflows/deploy.yml
name: Deploy Network Changes
on:
  push:
    branches: [main]
jobs:
  deploy:
    runs-on: ubuntu-latest
    steps:
      - uses: actions/checkout@v2
      - name: Run Ansible Playbook
        run: |
          ansible-playbook -i inventory \\
            playbooks/configure.yml"""
    
    text_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = workflow_example
    p.font.name = 'Courier New'
    p.font.size = Pt(14)
    
    # Slide 12: Implementation Phases
    create_content_slide(prs, "Implementation Roadmap", [
        "Phase 1: Planning & Assessment (2-3 weeks)",
        "Phase 2: Infrastructure Setup (2 weeks)",
        "Phase 3: Development & Testing (5 weeks)",
        "Phase 4: Training & Pilot (3-4 weeks)",
        "Phase 5: Production Rollout (4-6 weeks)",
        "Phase 6: Optimization & Closure (2 weeks)",
        "",
        "Total Timeline: ~6 months with parallel operations"
    ])
    
    # Slide 13: Success Metrics
    create_content_slide(prs, "Success Metrics", [
        "📉 Reduce change implementation time by 70%",
        "📉 Decrease configuration errors by 90%",
        "📈 Increase number of changes deployed per week by 300%",
        "✅ 100% of changes tracked in version control",
        "✅ 100% peer review compliance",
        "📊 Mean time to rollback < 5 minutes",
        "🎓 100% team trained and certified"
    ])
    
    # Slide 14: Risk Mitigation
    create_content_slide(prs, "Risk Mitigation Strategy", [
        "✓ Comprehensive lab testing environment",
        "✓ Phased rollout starting with non-critical devices",
        "✓ Parallel run with legacy methods during transition",
        "✓ Extensive training and documentation",
        "✓ Rollback procedures tested and documented",
        "✓ Change advisory board review for initial deployments",
        "✓ 24/7 support during critical deployment phases"
    ])
    
    # Slide 15: Team Skills & Training
    create_content_slide(prs, "Required Skills & Training", [
        "Git/GitHub fundamentals and workflows",
        "YAML syntax and best practices",
        "Ansible basics and playbook development",
        "Cisco NaC Collections modules",
        "GitHub Actions workflow creation",
        "Infrastructure as Code principles",
        "DevOps culture and collaboration",
        "",
        "Training: 2-week intensive + ongoing mentorship"
    ])
    
    # Slide 16: Security & Compliance
    create_content_slide(prs, "Security & Compliance", [
        "🔐 Credentials stored in GitHub Secrets (encrypted)",
        "👥 Role-based access control (RBAC)",
        "📝 All changes logged with author, timestamp, reason",
        "✅ Mandatory peer review before deployment",
        "🔍 Automated compliance checks in CI/CD pipeline",
        "📊 Integration with existing ITSM/ticketing systems",
        "🛡️ Network device access limited to Ansible control node"
    ])
    
    # Slide 17: Next Steps
    create_content_slide(prs, "Next Steps", [
        "1. Stakeholder alignment and budget approval",
        "2. Form project team (engineers, DevOps, security)",
        "3. Conduct current state assessment",
        "4. Set up lab environment for proof of concept",
        "5. Develop training curriculum",
        "6. Create detailed project plan and timeline",
        "7. Begin Phase 1: Planning & Assessment"
    ])
    
    # Slide 18: Q&A
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    title_frame = title_box.text_frame
    title_frame.text = "Questions & Discussion"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(54)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # Save presentation
    prs.save('C:\\Users\\danda\\NaC_Transformation_Presentation.pptx')
    print("PowerPoint presentation saved as: C:\\Users\\danda\\NaC_Transformation_Presentation.pptx")

if __name__ == "__main__":
    main()
