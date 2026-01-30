# Create a PowerPoint-ready single-slide Gantt chart for a 9-month Network as Code transformation

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Initialize presentation
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Title slide layout
slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.7))
title_tf = title_box.text_frame
title_tf.clear()
p = title_tf.paragraphs[0]
p.text = "Network as Code (NaC) Transformation — 9-Month Gantt"
p.font.size = Pt(28)
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

# Subtitle
sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.4))
sub_tf = sub_box.text_frame
sub_tf.clear()
p = sub_tf.paragraphs[0]
p.text = "From Manual CLI/GUI Changes to CI/CD with GitHub & Ansible Automation Platform"
p.font.size = Pt(14)
p.alignment = PP_ALIGN.CENTER

# Timeline settings
left_margin = Inches(1.0)
top_start = Inches(1.7)
row_height = Inches(0.45)
timeline_width = Inches(11.5)
months = 9
month_width = timeline_width / months

# Month headers
for i in range(months):
    box = slide.shapes.add_textbox(left_margin + month_width * i, top_start, month_width, Inches(0.4))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = f"Month {i+1}"
    p.font.size = Pt(11)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

# Gantt tasks: (Name, start_month (1-based), duration_months, color)
tasks = [
    ("Strategy & Buy-in", 1, 2, RGBColor(91, 155, 213)),
    ("Platform Setup", 2, 2, RGBColor(237, 125, 49)),
    ("Repo Standards & NaC Design", 3, 2, RGBColor(165, 165, 165)),
    ("Pilot Automation (Non-Prod)", 4, 2, RGBColor(112, 173, 71)),
    ("CI/CD Integration", 5, 2, RGBColor(68, 114, 196)),
    ("Production Readiness", 6, 2, RGBColor(255, 192, 0)),
    ("Initial Production Rollout", 7, 2, RGBColor(84, 130, 53)),
    ("Expansion & Standardization", 8, 2, RGBColor(191, 143, 0)),
    ("Operationalization", 9, 1, RGBColor(128, 128, 128)),
]

# Draw task labels and bars
for idx, (name, start, duration, color) in enumerate(tasks):
    y = top_start + Inches(0.5) + row_height * idx

    # Task label
    label = slide.shapes.add_textbox(Inches(0.3), y, Inches(0.65), Inches(0.35))
    tf = label.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(11)

    # Task bar
    bar_left = left_margin + month_width * (start - 1)
    bar_width = month_width * duration
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        bar_left,
        y,
        bar_width,
        Inches(0.3)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

# Save presentation
file_path = "/mnt/data/NaC_9_Month_Gantt.pptx"
prs.save(file_path)

file_path
